from fastapi import FastAPI, Form, File, UploadFile, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageOps
from sqlalchemy import create_engine, MetaData, Table, Column, String, Integer, Text, DateTime
from sqlalchemy.sql import func
import databases
import uuid

#DATABASE_URL = "postgresql://postgres:Sevana%401995@db.dnwfdswyuqahzmkgewui.supabase.co:5432/postgres"
DATABASE_URL = "postgresql://football_db_9sfj_user:O2Ux8ljm3l1XCkJjT2hwcS4m6dTfRSO5@dpg-d1laqcadbo4c739pggn0-a/football_db_9sfj"
database = databases.Database(DATABASE_URL)
metadata = MetaData()

players = Table(
    "players",
    metadata,
    Column("id", String, primary_key=True),
    Column("name", Text),
    Column("house", Text),
    Column("mobile", Text),
    Column("whatsapp", Text),
    Column("father", Text),
    Column("age", Integer),
    Column("unit", Text),
    Column("epl", Text),
    Column("prev_team", Text),
    Column("photo_path", Text),
    Column("registered_at", DateTime, server_default=func.now())
)

engine = create_engine(DATABASE_URL)
metadata.create_all(engine)

app = FastAPI()

@app.on_event("startup")
async def startup():
    await database.connect()

@app.on_event("shutdown")
async def shutdown():
    await database.disconnect()

# CORS setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Paths
UPLOAD_DIR = "backend/uploads"
LOGO_PATH = "backend/logo.png"

# Ensure upload folder exists
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Static mounts
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")
app.mount("/static", StaticFiles(directory="frontend"), name="static")
app.mount("/images", StaticFiles(directory="backend"), name="images")

@app.get("/", response_class=HTMLResponse)
def root():
    return FileResponse("frontend/index.html")

@app.get("/register", response_class=HTMLResponse)
def serve_register():
    return FileResponse("frontend/register.html")

@app.post("/register")
async def register(
    name: str = Form(...),
    house: str = Form(...),
    mobile: str = Form(...),
    whatsapp: str = Form(...),
    father: str = Form(...),
    age: int = Form(...),
    unit: str = Form(...),
    epl: str = Form(...),
    prev_team: str = Form(""),
    photo: UploadFile = File(...)
):
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    safe_name = name.replace(" ", "_")
    filename = f"{safe_name}_{timestamp}.jpg"
    file_path = os.path.join(UPLOAD_DIR, filename)

    # Save photo
    with open(file_path, "wb") as f:
        content = await photo.read()
        f.write(content)

    relative_path = f"uploads/{filename}"

    player_data = {
        "id": str(uuid.uuid4()),
        "name": name,
        "house": house,
        "mobile": mobile,
        "whatsapp": whatsapp,
        "father": father,
        "age": age,
        "unit": unit,
        "epl": epl,
        "prev_team": prev_team,
        "photo_path": relative_path
    }

    query = players.insert().values(**player_data)
    await database.execute(query)

    return JSONResponse({"message": "✅ Registration successful!"})

@app.get("/admin-data")
async def get_admin_data():
    query = players.select()
    result = await database.fetch_all(query)
    return [dict(row) for row in result]

@app.get("/download-ppt")
def download_ppt():
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    with engine.connect() as conn:
        result = conn.execute(players.select()).fetchall()

    for idx, row in enumerate(result):
        slide = prs.slides.add_slide(blank_layout)

        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3), Inches(9), Inches(6.5)
        )
        card.fill.solid()
        bg_color = RGBColor(235, 245, 255) if idx % 2 == 0 else RGBColor(250, 250, 250)
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = RGBColor(180, 180, 180)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = row.name
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 70, 140)

        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.5))
        tf = info_box.text_frame
        tf.clear()
        fields = [
            ("Unit", row.unit),
            ("House", row.house),
            ("Father's Name", row.father),
            ("Age", row.age),
            ("Previous Team", row.prev_team or "N/A"),
            ("Base Value", "30")
        ]
        for label, value in fields:
            para = tf.add_paragraph()
            para.text = f"{label}: {value}"
            para.font.size = Pt(22)
            para.font.bold = True
            para.font.color.rgb = RGBColor(30, 30, 30)

        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(7.5), Inches(0.3), width=Inches(1.5))

        footer_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.6), Inches(9), Inches(0.5))
        para = footer_box.text_frame.paragraphs[0]
        para.text = "EPL Tournament 2025"
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(100, 100, 100)
        para.alignment = PP_ALIGN.CENTER

        photo_path = os.path.join("backend", row.photo_path)
        if os.path.exists(photo_path):
            with Image.open(photo_path) as img:
                bordered = ImageOps.expand(img, border=5, fill='black')
                temp_path = f"backend/_temp_{idx}.jpg"
                bordered.save(temp_path)
            slide.shapes.add_picture(temp_path, Inches(6.2), Inches(1.8), width=Inches(2.5), height=Inches(2.5))

    output_file = "backend/player_data.pptx"
    prs.save(output_file)

    return FileResponse(output_file, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="players.pptx")

@app.get("/admin", response_class=HTMLResponse)
def serve_admin():
    return FileResponse("frontend/admin.html")

import io
import csv
from fastapi.responses import StreamingResponse

@app.get("/download-csv")
async def download_csv():
    query = players.select()
    rows = await database.fetch_all(query)

    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow([
        "name", "house", "mobile", "whatsapp", "father", "age",
        "unit", "epl", "prev_team", "photo_path", "registered_at"
    ])
    for row in rows:
        writer.writerow([
            row["name"], row["house"], row["mobile"], row["whatsapp"],
            row["father"], row["age"], row["unit"], row["epl"],
            row["prev_team"], row["photo_path"], row["registered_at"]
        ])
    output.seek(0)
    return StreamingResponse(output, media_type="text/csv", headers={"Content-Disposition": "attachment; filename=players.csv"})
